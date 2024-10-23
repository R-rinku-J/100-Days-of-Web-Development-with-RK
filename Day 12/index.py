from pptx import Presentation
from pptx.util import Inches
import requests

# Create a PowerPoint presentation object
ppt = Presentation()

# Slide 1: Title Slide
slide_layout = ppt.slide_layouts[0]  # 0 is the layout for the title slide
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Natural Disaster Maintenance Model"
subtitle.text = "Real-Time Image Validation and Relief Delivery System"

# Slide 2: Introduction Slide
slide_layout = ppt.slide_layouts[1]  # 1 is the layout for a content slide
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Introduction"

content = slide.shapes.placeholders[1].text_frame
content.text = (
    "The project focuses on building a model for natural disaster maintenance "
    "where an image is uploaded to the system. The model validates in real-time "
    "if the image was taken within a specific period and helps prioritize "
    "relief delivery to the most affected areas."
)

# Adding an image to the introduction slide
image_url = "https://i.imgur.com/GOzO2mO.png"  # Disaster scene image
image_path = "intro_image.png"
img_data = requests.get(image_url).content
with open(image_path, 'wb') as handler:
    handler.write(img_data)
slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), Inches(4), Inches(3))

# Slide 3: Objective
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Objective"

content = slide.shapes.placeholders[1].text_frame
content.text = (
    "The objective is to create a system that can validate the time of a "
    "submitted image, ensuring it's within a recent timeframe. Upon validation, "
    "the system helps provide rapid and prioritized relief to disaster-affected regions."
)

# Adding an image to the objective slide
image_url = "https://i.imgur.com/AmRw77K.png"  # Objective illustration
image_path = "objective_image.png"
img_data = requests.get(image_url).content
with open(image_path, 'wb') as handler:
    handler.write(img_data)
slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), Inches(4), Inches(3))

# Slide 4: Features
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Features"

content = slide.shapes.placeholders[1].text_frame
content.text = (
    "1. Real-time image validation for authenticity.\n"
    "2. Rapid response to the most affected areas based on image data.\n"
    "3. Seamless coordination with relief teams.\n"
    "4. System prioritization based on severity."
)

# Adding an image to the features slide
image_url = "https://i.imgur.com/vwRjTsn.png"  # Features illustration
image_path = "features_image.png"
img_data = requests.get(image_url).content
with open(image_path, 'wb') as handler:
    handler.write(img_data)
slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), Inches(4), Inches(3))

# Slide 5: Technologies Used
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Technologies Used"

content = slide.shapes.placeholders[1].text_frame
content.text = (
    "1. Machine Learning: For validating image authenticity and identifying affected areas.\n"
    "2. Python & Flask: Backend development and API integration.\n"
    "3. HTML, CSS, JavaScript: Frontend for user interface and image uploads.\n"
    "4. Cloud Storage: For storing validated images and records."
)

# Adding an image to the technologies slide
image_url = "https://i.imgur.com/jv8i67J.png"  # Technologies diagram
image_path = "technologies_image.png"
img_data = requests.get(image_url).content
with open(image_path, 'wb') as handler:
    handler.write(img_data)
slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), Inches(4), Inches(3))

# Slide 6: Conclusion
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Conclusion"

content = slide.shapes.placeholders[1].text_frame
content.text = (
    "This model can significantly improve the response time and effectiveness "
    "of disaster relief efforts. By ensuring that resources are directed to the "
    "most affected areas in real-time, we can minimize the damage and provide "
    "support where it is needed the most."
)

# Adding an image to the conclusion slide
image_url = "https://i.imgur.com/lodurHV.png"  # Rescue operations image
image_path = "conclusion_image.png"
img_data = requests.get(image_url).content
with open(image_path, 'wb') as handler:
    handler.write(img_data)
slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), Inches(4), Inches(3))

# Save the presentation
ppt_path = "Natural_Disaster_Maintenance_Model_Presentation_with_Images.pptx"
ppt.save(ppt_path)

print(f"Presentation saved as {ppt_path}")
